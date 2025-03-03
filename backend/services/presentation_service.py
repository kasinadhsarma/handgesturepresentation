import os
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation
from PIL import Image
import io
from typing import List, Dict
from config.settings import settings

class PresentationService:
    def __init__(self):
        self.upload_dir = settings.UPLOAD_DIR
        os.makedirs(self.upload_dir, exist_ok=True)
        
    async def process_presentation(self, file: bytes, filename: str) -> Dict:
        """Process uploaded presentation file (PPT/PPTX/PDF)."""
        file_ext = filename.split('.')[-1].lower()
        
        if file_ext in ['ppt', 'pptx']:
            return await self._process_powerpoint(file)
        elif file_ext == 'pdf':
            return await self._process_pdf(file)
        else:
            raise ValueError("Unsupported file format")
    
    async def _process_powerpoint(self, file: bytes) -> Dict:
        """Convert PowerPoint to images and extract text."""
        presentation = Presentation(io.BytesIO(file))
        slides = []
        
        for i, slide in enumerate(presentation.slides):
            # Save slide as image
            img_path = f"{self.upload_dir}/slide_{i+1}.png"
            self._save_slide_as_image(slide, img_path)
            
            # Extract text content
            text_content = self._extract_slide_text(slide)
            
            slides.append({
                "number": i + 1,
                "image_path": img_path,
                "text_content": text_content
            })
        
        return {
            "total_slides": len(slides),
            "slides": slides
        }
    
    async def _process_pdf(self, file: bytes) -> Dict:
        """Convert PDF to images and extract text."""
        pdf_document = fitz.open(stream=file, filetype="pdf")
        slides = []
        
        for i in range(pdf_document.page_count):
            page = pdf_document[i]
            
            # Convert page to image
            pix = page.get_pixmap()
            img_path = f"{self.upload_dir}/slide_{i+1}.png"
            pix.save(img_path)
            
            # Extract text
            text_content = page.get_text()
            
            slides.append({
                "number": i + 1,
                "image_path": img_path,
                "text_content": text_content
            })
        
        return {
            "total_slides": len(slides),
            "slides": slides
        }
    
    def _save_slide_as_image(self, slide, path: str):
        """Save PowerPoint slide as image."""
        # Implementation depends on the specific library used
        # This is a placeholder for the actual implementation
        pass
    
    def _extract_slide_text(self, slide) -> str:
        """Extract text content from PowerPoint slide."""
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        
        return "\n".join(text_content)
    
    async def save_annotation(
        self,
        presentation_id: str,
        slide_number: int,
        annotation_data: str
    ) -> Dict:
        """Save slide annotation."""
        # Convert base64 image data to PIL Image
        img_data = annotation_data.split(',')[1]
        img = Image.open(io.BytesIO(bytes(img_data, 'utf-8')))
        
        # Save annotated slide
        annotation_path = f"{self.upload_dir}/{presentation_id}_slide_{slide_number}_annotated.png"
        img.save(annotation_path)
        
        return {
            "success": True,
            "annotation_path": annotation_path
        }

